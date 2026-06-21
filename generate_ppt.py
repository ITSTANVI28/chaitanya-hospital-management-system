import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Define color scheme (Medical Deep Blue, Light Slate, Crimson Red, Teal)
    DARK_BLUE = RGBColor(12, 35, 64)       # Primary dark color
    TEAL = RGBColor(0, 150, 136)           # Secondary accent
    CHARCOAL = RGBColor(51, 51, 51)        # Body text
    LIGHT_GRAY = RGBColor(245, 247, 250)   # Slide background
    WHITE = RGBColor(255, 255, 255)
    RED = RGBColor(211, 47, 47)
    
    # Slide Dimensions (16:9 widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Helper to set slide background
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to create professional text boxes
    def add_textbox(slide, left, top, width, height):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        return tf

    # Helper to format slide title
    def add_slide_header(slide, title_text):
        tf = add_textbox(slide, Inches(0.75), Inches(0.4), Inches(11.833), Inches(0.8))
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Poppins'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        # Add a subtle teal underline for visual elegance
        shape = slide.shapes.add_shape(
            1, # Rectangle
            Inches(0.75), Inches(1.15), Inches(2.0), Inches(0.06)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = TEAL
        shape.line.fill.background() # No line

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Background)
    # -------------------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    set_background(slide1, DARK_BLUE)
    
    # Accent shape on left edge
    accent = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    
    tf1 = add_textbox(slide1, Inches(1.5), Inches(2.2), Inches(10.5), Inches(3.0))
    p1 = tf1.paragraphs[0]
    p1.text = "Chaitanya Multi Speciality Hospital"
    p1.font.name = 'Poppins'
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    
    p2 = tf1.add_paragraph()
    p2.text = "Hospital Management System (HMS) — Working Report"
    p2.font.name = 'Poppins'
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEAL
    p2.space_before = Pt(10)
    
    p3 = tf1.add_paragraph()
    p3.text = "Prepared by: Pranav (Intern)\nDate: June 2026"
    p3.font.name = 'Poppins'
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(180, 190, 210)
    p3.space_before = Pt(40)

    # -------------------------------------------------------------
    # SLIDE 2: Project Overview & Context (Light Background)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_background(slide2, LIGHT_GRAY)
    add_slide_header(slide2, "Project Overview & Context")
    
    tf2 = add_textbox(slide2, Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    
    p = tf2.paragraphs[0]
    p.text = "From Static Design to Full-Stack Application"
    p.font.name = 'Poppins'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(14)
    
    bullets = [
        "**Original Internship Plan:** Build a simple, static front-end website with simulated, in-memory data that resets on page refresh.",
        "**Advanced Implementation:** Upgraded the system to a database-driven web application with a PHP backend API and a persistent MySQL database.",
        "**User Interface Aesthetics:** Implemented modern web standards with Google Fonts (Poppins), custom variables, smooth transitions, and glassmorphism elements.",
        "**Dual Interfaces:** Public-facing information & appointment booking portal (`index.html`) integrated with an Administrator Dashboard (`dashboard.html`)."
    ]
    
    for bullet in bullets:
        p = tf2.add_paragraph()
        p.space_before = Pt(12)
        p.font.name = 'Poppins'
        p.font.size = Pt(16)
        p.font.color.rgb = CHARCOAL
        
        # Simple markdown bold parser
        if "**" in bullet:
            parts = bullet.split("**")
            run1 = p.add_run()
            run1.text = "• " + parts[1]
            run1.font.bold = True
            run1.font.color.rgb = DARK_BLUE
            
            run2 = p.add_run()
            run2.text = parts[2]
        else:
            p.text = "• " + bullet

    # -------------------------------------------------------------
    # SLIDE 3: System Architecture (Three-Column Layout)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    set_background(slide3, LIGHT_GRAY)
    add_slide_header(slide3, "System Architecture")
    
    # Column 1: Frontend
    tf_col1 = add_textbox(slide3, Inches(0.75), Inches(1.8), Inches(3.6), Inches(4.8))
    p_col1 = tf_col1.paragraphs[0]
    p_col1.text = "1. Frontend (UI/UX)"
    p_col1.font.bold = True
    p_col1.font.size = Pt(18)
    p_col1.font.color.rgb = DARK_BLUE
    p_col1.space_after = Pt(10)
    
    col1_points = [
        "**Pure HTML5 & CSS3:** No bloated frameworks; clean, hand-written styling sheets.",
        "**Vanilla JavaScript:** Handles DOM updates, modal popups, and AJAX requests dynamically.",
        "**Interactive Modules:** Filterable doctor profiles, dynamic inpatient bed layouts."
    ]
    for pt in col1_points:
        p = tf_col1.add_paragraph()
        p.font.size = Pt(14)
        p.font.color.rgb = CHARCOAL
        p.space_before = Pt(8)
        parts = pt.split("**")
        run1 = p.add_run()
        run1.text = "• " + parts[1]
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = parts[2]

    # Column 2: Backend API
    tf_col2 = add_textbox(slide3, Inches(4.85), Inches(1.8), Inches(3.6), Inches(4.8))
    p_col2 = tf_col2.paragraphs[0]
    p_col2.text = "2. Backend API Layer"
    p_col2.font.bold = True
    p_col2.font.size = Pt(18)
    p_col2.font.color.rgb = DARK_BLUE
    p_col2.space_after = Pt(10)
    
    col2_points = [
        "**Modular PHP Scripts:** API endpoints hosted under `/api` for decoupled, quick request processing.",
        "**PDO Interface:** Safe query execution with database exception handling.",
        "**Direct Communication:** Seamless fetching of JSON arrays via standard HTTP methods."
    ]
    for pt in col2_points:
        p = tf_col2.add_paragraph()
        p.font.size = Pt(14)
        p.font.color.rgb = CHARCOAL
        p.space_before = Pt(8)
        parts = pt.split("**")
        run1 = p.add_run()
        run1.text = "• " + parts[1]
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = parts[2]

    # Column 3: Database
    tf_col3 = add_textbox(slide3, Inches(8.95), Inches(1.8), Inches(3.6), Inches(4.8))
    p_col3 = tf_col3.paragraphs[0]
    p_col3.text = "3. MySQL Database"
    p_col3.font.bold = True
    p_col3.font.size = Pt(18)
    p_col3.font.color.rgb = DARK_BLUE
    p_col3.space_after = Pt(10)
    
    col3_points = [
        "**Schema `chaitanya_hms`:** Houses 6 dedicated tables mapping hospital assets.",
        "**Referential Integrity:** Structured ENUM fields for tracking stages (e.g., Patient Status, Shift Timings).",
        "**In-Transaction Commits:** Transaction locks used during stock processing to ensure data consistency."
    ]
    for pt in col3_points:
        p = tf_col3.add_paragraph()
        p.font.size = Pt(14)
        p.font.color.rgb = CHARCOAL
        p.space_before = Pt(8)
        parts = pt.split("**")
        run1 = p.add_run()
        run1.text = "• " + parts[1]
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = parts[2]

    # -------------------------------------------------------------
    # SLIDE 4: Key Modules - Patient & Bed Management
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_background(slide4, LIGHT_GRAY)
    add_slide_header(slide4, "Patient & Bed Management")
    
    tf4 = add_textbox(slide4, Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    p = tf4.paragraphs[0]
    p.text = "Optimized Patient Admission & Bed Allocation"
    p.font.name = 'Poppins'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(14)
    
    points4 = [
        "**Dynamic Bed Visualizer:** Real-time representation of Ward Beds (GEN, ICU, PED, MAT, ORT). Occupied beds highlight in Red; vacant beds highlight in Green.",
        "**Interactive Inpatient Records:** Live searching by patient name or bed ID and filtering by ward department.",
        "**Status Tracking:** Color-coded badges for active tracking: Stable (Green), Critical (Red), Post-Op (Orange), Discharge (Blue).",
        "**Operational Actions:** Seamless patient admissions with validation, instant discharges (which frees up the occupied bed in the layout), and record deletion."
    ]
    for pt in points4:
        p = tf4.add_paragraph()
        p.font.size = Pt(16)
        p.font.color.rgb = CHARCOAL
        p.space_before = Pt(12)
        parts = pt.split("**")
        run1 = p.add_run()
        run1.text = "• " + parts[1]
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = parts[2]

    # -------------------------------------------------------------
    # SLIDE 5: Key Modules - Inventory & Barcode Scanner
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_background(slide5, LIGHT_GRAY)
    add_slide_header(slide5, "Stock Inventory & Barcode Scanning")
    
    tf5 = add_textbox(slide5, Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    p = tf5.paragraphs[0]
    p.text = "Intelligent Medicine & Consumables Tracking"
    p.font.name = 'Poppins'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(14)
    
    points5 = [
        "**Low-Stock Alerting:** Live indicators pointing out items falling below their specified minimum threshold (marked with a prominent red background and 'Restock Urgent' badge).",
        "**Barcode/QR Scanner Integration:** Simulated scanning modal supporting Live Camera (via html5-qrcode), photo capture/upload, and manual entry.",
        "**Dynamic Form Autofill:** Scanning an EAN-13 code identifies the item name and type, instantly autofilling the restock form details.",
        "**Stock Orders Log:** Keeps records of all pending orders. Clicking 'Receive Order' triggers a backend transaction that increments the inventory count automatically."
    ]
    for pt in points5:
        p = tf5.add_paragraph()
        p.font.size = Pt(16)
        p.font.color.rgb = CHARCOAL
        p.space_before = Pt(12)
        parts = pt.split("**")
        run1 = p.add_run()
        run1.text = "• " + parts[1]
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = parts[2]

    # -------------------------------------------------------------
    # SLIDE 6: Key Modules - Staff & Waste Management
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_background(slide6, LIGHT_GRAY)
    add_slide_header(slide6, "Staff Info & Waste Management")
    
    tf6 = add_textbox(slide6, Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    p = tf6.paragraphs[0]
    p.text = "Hospital Shift Rosters & Bio-Hazard Disposal Control"
    p.font.name = 'Poppins'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(14)
    
    points6 = [
        "**Shift Roster:** Organizes staff into Morning, Evening, and Night shifts with statuses: 'On Duty', 'Off Duty', and 'On Leave'.",
        "**Waste Log:** Logs weights of Biomedical (Red), Sharps (Orange), and General (Gray) waste to maintain compliance.",
        "**Dispose Cart:** An automated routing system for failed waste entries. If disposal fails, the entry goes to the Dispose Cart.",
        "**Actionable Failures:** From the Dispose Cart, admins can click 'Retry' to resubmit the log as Pending or permanently delete it."
    ]
    for pt in points6:
        p = tf6.add_paragraph()
        p.font.size = Pt(16)
        p.font.color.rgb = CHARCOAL
        p.space_before = Pt(12)
        parts = pt.split("**")
        run1 = p.add_run()
        run1.text = "• " + parts[1]
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = parts[2]

    # -------------------------------------------------------------
    # SLIDE 7: Recent Technical Upgrades
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(slide_layout)
    set_background(slide7, LIGHT_GRAY)
    add_slide_header(slide7, "Recent Technical Upgrades")
    
    tf7 = add_textbox(slide7, Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    p = tf7.paragraphs[0]
    p.text = "Key Upgrades Implemented"
    p.font.name = 'Poppins'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(14)
    
    points7 = [
        "**Automated Dispose Cart:** Implemented `api/retry_waste.php` and `api/delete_waste_log.php` to handle failed waste records dynamically.",
        "**Dynamic Inventory Auto-Increment:** Connected the 'Receive Order' button in the dashboard to `api/receive_stock_order.php` for seamless transactions.",
        "**Bug Fixes & Refactoring:** Resolved the Staff department loading issue and implemented robust, error-checked SQL schema scripts.",
        "**Data Consistency:** Fully migrated the static mock arrays to MySQL database tables, making the dashboard fully operational."
    ]
    for pt in points7:
        p = tf7.add_paragraph()
        p.font.size = Pt(16)
        p.font.color.rgb = CHARCOAL
        p.space_before = Pt(12)
        parts = pt.split("**")
        run1 = p.add_run()
        run1.text = "• " + parts[1]
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = parts[2]

    # -------------------------------------------------------------
    # SLIDE 8: Project Highlights & Conclusion (Dark Background)
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(slide_layout)
    set_background(slide8, DARK_BLUE)
    
    # Accent shape on left edge
    accent8 = slide8.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    accent8.fill.solid()
    accent8.fill.fore_color.rgb = TEAL
    accent8.line.fill.background()
    
    tf8 = add_textbox(slide8, Inches(1.5), Inches(1.8), Inches(10.5), Inches(4.5))
    p8_1 = tf8.paragraphs[0]
    p8_1.text = "Project Highlights & Conclusion"
    p8_1.font.name = 'Poppins'
    p8_1.font.size = Pt(36)
    p8_1.font.bold = True
    p8_1.font.color.rgb = WHITE
    p8_1.space_after = Pt(20)
    
    conclusions = [
        "Modern, visual, and highly responsive UI/UX designed using standard CSS3 styling variables.",
        "Robust database integration that implements real-world medical workflows like bed allocation, low-stock warnings, and barcode scanner entries.",
        "A highly scalable codebase suitable for further feature additions (e.g., patient billing, laboratory records).",
        "Demonstrates complete full-stack capabilities using foundational, clean web development stacks."
    ]
    
    for c in conclusions:
        p = tf8.add_paragraph()
        p.font.name = 'Poppins'
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(220, 225, 235)
        p.space_before = Pt(12)
        p.text = "✔  " + c
        
    prs.save("Chaitanya_Hospital_HMS_Working_Report.pptx")
    print("Presentation generated successfully as Chaitanya_Hospital_HMS_Working_Report.pptx")

if __name__ == "__main__":
    create_presentation()
